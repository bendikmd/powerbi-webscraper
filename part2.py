from openpyxl import load_workbook
from pptx import Presentation

# Load the Excel file
wb = load_workbook('powerbi_blog_articles.xlsx')
ws = wb.active

# Load the PowerPoint template
prs = Presentation('template1.pptx')

# Iterate through each row in the Excel file, skipping the header row
for row in ws.iter_rows(min_row=2, values_only=True):
    # Create a new slide based on the desired layout
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)

    # Map placeholders to column names
    placeholders = {
        shape.name: value
        for shape, value in zip(slide.shapes, row)
        if shape.has_text_frame
    }

    # Populate placeholders with the corresponding data
    for shape in slide.shapes:
        if shape.name in placeholders:
            shape.text_frame.text = str(placeholders[shape.name])

# Save the modified PowerPoint presentation with a custom name
output_presentation = 'generated_presentation.pptx'
prs.save(output_presentation)

# Print a success message
print(f"Process completed successfully. Generated presentation saved as '{output_presentation}'")
